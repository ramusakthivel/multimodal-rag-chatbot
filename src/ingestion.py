import os
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

def extract_text_from_pdf(file_path):
    """Extracts raw text page-by-page from a PDF file."""
    text = ""
    try:
        reader = PdfReader(file_path)
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += f"\n--- Page {page_num + 1} ---\n{page_text}"
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
    return text

def extract_text_from_docx(file_path):
    """Extracts raw text from paragraphs and simple tables in a Word document."""
    text = ""
    try:
        doc = Document(file_path)
        # Extract text from standard paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        
        # Extract text from simple text tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text += " | ".join(row_text) + "\n"
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path):
    """Extracts raw text from shapes and notes inside a PowerPoint presentation."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
    return text

def process_raw_documents(directory_path="data/raw"):
    """Scans the directory and runs the appropriate parser based on file extension."""
    documents_data = {}
    
    if not os.path.exists(directory_path):
        print(f"Directory '{directory_path}' not found. Creating it now...")
        os.makedirs(directory_path)
        return documents_data

    for filename in os.listdir(directory_path):
        file_path = os.path.join(directory_path, filename)
        ext = filename.split('.')[-1].lower()
        
        print(f"Processing: {filename}...")
        
        if ext == 'pdf':
            documents_data[filename] = extract_text_from_pdf(file_path)
        elif ext == 'docx':
            documents_data[filename] = extract_text_from_docx(file_path)
        elif ext == 'pptx':
            documents_data[filename] = extract_text_from_pptx(file_path)
        else:
            print(f"Skipping unsupported file type: {filename}")
            
    return documents_data

if __name__ == "__main__":
    # Quick local test of our ingestion pipeline
    print("Testing Ingestion Service...")
    extracted_data = process_raw_documents()
    print(f"Successfully processed {len(extracted_data)} files.")